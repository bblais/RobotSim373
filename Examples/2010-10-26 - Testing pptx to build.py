#!/usr/bin/env python
# coding: utf-8

# In[1]:


get_ipython().run_line_magic('pylab', 'inline')


# In[2]:


from RobotSim373 import *


# In[3]:


from RobotSim373.utils import pptx2build


# In[4]:


pptx2build('test drawing.pptx')


# In[5]:


def build(robot):
    box1=Box(robot,x=3.25,y=17.36,angle=0.00,width=1.25,height=3.53,name="box1")
    box2=Box(robot,x=7.94,y=10.83,angle=326.93,width=1.15,height=4.18,name="box2")
    disk3=Disk(robot,x=5.02,y=12.30,angle=0.00,radius=1.77,name="disk3")
    disk4=Disk(robot,x=11.19,y=13.14,angle=0.00,radius=1.04,name="disk4")

    connect(disk3,box1,'weld')
    connect(disk3,box2,'weld')
    
    robot.offset(dx=5)


# In[6]:


def act(t,robot):
    robot['box1'].F=.1
    pass


# In[7]:


env=Environment(24,24)
robot=Robot(env)
build(robot)


# In[12]:


run_sim(env,act,
        total_time=1000,  # seconds
        dt=1/60,
        dt_display=5,  # make this larger for a faster display
       )


# In[5]:


fname='/Users/bblais/Downloads/Robot Drawing.pptx'
from pptx import Presentation

prs=Presentation(fname)


# In[6]:


slide=prs.slides[0]


# In[7]:


for shape in slide.shapes:
    print(shape.name)


# In[8]:


from RobotSim373.utils import pptx2build
pptx2build(fname)


# In[ ]:




