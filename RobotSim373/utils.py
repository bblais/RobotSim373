from pptx import Presentation

def convert(shape_info,env):
    from math import radians,cos
    scale=24/env['height']
        
    shape_info['width']*=scale
    shape_info['height']*=scale
    
    shape_info['left']-=env['left']
    shape_info['left']*=scale
    
    shape_info['top']-=env['top']
    shape_info['top']*=scale
    
    shape_info['rotation']=(360-shape_info['rotation'])%360
    
    
    if shape_info['type']=='Straight':
        shape_info['x1']-=env['left']
        shape_info['x1']*=scale
        shape_info['x2']-=env['left']
        shape_info['x2']*=scale

        shape_info['y1']-=env['top']
        shape_info['y1']*=scale
        shape_info['y1']=24-shape_info['y1']
        shape_info['y2']-=env['top']
        shape_info['y2']*=scale
        shape_info['y2']=24-shape_info['y2']
        
        
        shape_info['x']=shape_info['x1']
        shape_info['y']=shape_info['y1']
        
    else:
        shape_info['x']=shape_info['left']+shape_info['width']/2
        shape_info['y']=24-shape_info['top']
    
    

    return shape_info
    

def pptx2build(fname):

    prs=Presentation(fname)
    slide=prs.slides[0]

    env={}
    env['left']=0
    env['top']=0
    env['width']=prs.slide_width
    env['height']=prs.slide_height


    shapes_info=[]
    for shape in slide.shapes:
        shape_info={}
        shape_type=shape.name.split()[0]
        
        if shape_type=="Rounded" or 'Rounded' in shape.name:
            env['left']=shape.left
            env['top']=shape.top
            env['width']=shape.width
            env['height']=shape.height
            continue
        
        assert shape_type in ['Rectangle','Oval','Straight'], f"Only Rectangle and Circle Shapes Allowed: Found {shape_type}"
        
        if shape_type=='Oval':
            assert shape.width==shape.height, "Only Circle (no Oval) Shapes Allowed"
            
            
        shape_info['type']=shape_type
        shape_info['left']=shape.left
        shape_info['top']=shape.top
        shape_info['width']=shape.width
        shape_info['height']=shape.height
        shape_info['rotation']=shape.rotation
        
        
        if shape_info['type']=='Straight':
            shape_info['x1']=shape.begin_x
            shape_info['x2']=shape.end_x
            shape_info['y1']=shape.begin_y
            shape_info['y2']=shape.end_y
        
        
        
        shapes_info.append(shape_info)
        
        
            
        #print(shape.name,shape.rotation)
        
        
    shapes_info=[convert(_,env) for _ in shapes_info]    
    #print(shapes_info)
            

    print("def build(robot):")
    for i,shape in enumerate(shapes_info):
        x,y,w,h,a=shape['x'],shape['y'],shape['width'],shape['height'],shape['rotation']
        
        
        if shape['type']=='Rectangle':
            print(f'    box{i+1}=Box(robot,x={x:0.2f},y={y:0.2f},angle={a:0.2f},width={w:0.2f},height={h:0.2f},name="box{i+1}")')
        elif shape['type']=='Oval':
            print(f'    disk{i+1}=Disk(robot,x={x:0.2f},y={y:0.2f},angle={a:0.2f},radius={w/2:0.2f},name="disk{i+1}")')
        elif shape['type']=='Straight':
            pass
        else:
            raise ValueError
            
    print()
            
    for i,shape in enumerate(shapes_info):
        if not shape['type']=='Straight':
            continue
            
        d1=1e500
        d2=1e500
        i1=-1
        i2=-1
        
        for j,shape2 in enumerate(shapes_info):
            if shape2['type']=='Straight':
                continue
            
            d=(shape['x1']-shape2['x'])**2+(shape['y1']-shape2['y'])**2
            if d<d1:
                d1=d
                i1=j
            d=(shape['x2']-shape2['x'])**2+(shape['y2']-shape2['y'])**2
            if d<d2:
                d2=d
                i2=j
            
            
        if shapes_info[i1]['type']=='Rectangle':
            v1=f"box{i1+1}"
        elif shapes_info[i1]['type']=='Oval':
            v1=f"disk{i1+1}"
        else:
            raise ValueError

        if shapes_info[i2]['type']=='Rectangle':
            v2=f"box{i2+1}"
        elif shapes_info[i2]['type']=='Oval':
            v2=f"disk{i2+1}"
        else:
            raise ValueError
            
            
        print(f"    connect({v1},{v2},'weld')")
                        